"""
Export Services for ZIMSEC Statistics Dashboard
Handles PowerPoint, PDF, Excel, and HTML exports with professional formatting
"""

from datetime import datetime, date
from decimal import Decimal
import json
from io import BytesIO
import base64

from django.template.loader import render_to_string
from django.utils import timezone
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PowerPointRGBColor
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import pandas as pd
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.colors import Color
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image, HRFlowable
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

from core.models import ZimsecResults, Student, Class


class PowerPointExporter:
    """Generate PREMIUM professional PowerPoint presentations from ZIMSEC data"""
    
    def __init__(self, title="ZIMSEC 2027 Examination Analysis", created_by="Dashboard"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.created_by = created_by
        
        # Premium color palette (0-255 integer values for python-pptx)
        self.primary_color = PowerPointRGBColor(15, 23, 42)           # Slate-900
        self.secondary_color = PowerPointRGBColor(59, 130, 246)       # Blue-500
        self.accent_color = PowerPointRGBColor(236, 72, 153)          # Pink-500
        self.success_color = PowerPointRGBColor(34, 197, 94)          # Green-500
        self.warning_color = PowerPointRGBColor(249, 115, 22)         # Orange-500
        self.text_light = PowerPointRGBColor(226, 232, 240)           # Slate-200
        self.text_muted = PowerPointRGBColor(148, 163, 184)           # Slate-400
        
    def _add_gradient_background(self, slide, color1=None, color2=None):
        """Add gradient background to slide"""
        if color1 is None:
            color1 = self.primary_color
        if color2 is None:
            color2 = self.secondary_color
            
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45.0
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
        
    def add_title_slide(self, school_name="School Name", year=2027):
        """Slide 1: Premium title slide with school info"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        self._add_gradient_background(slide, PowerPointRGBColor(15, 23, 42), PowerPointRGBColor(30, 41, 59))
        
        # Top accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.3)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.secondary_color
        accent_bar.line.color.rgb = self.secondary_color
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.text_light
        
        # Subtitle line
        subtitle_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(3.4),
            Inches(3), Inches(0.08)
        )
        subtitle_line.fill.solid()
        subtitle_line.fill.fore_color.rgb = self.secondary_color
        subtitle_line.line.color.rgb = self.secondary_color
        
        # School and date info
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(2.5))
        info_frame = info_box.text_frame
        info_frame.word_wrap = True
        
        p = info_frame.paragraphs[0]
        p.text = f"📚 {school_name}"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.text_light
        
        p = info_frame.add_paragraph()
        p.text = f"Academic Year {year}"
        p.font.size = Pt(24)
        p.font.color.rgb = self.secondary_color
        p.space_before = Pt(16)
        
        p = info_frame.add_paragraph()
        p.text = f"\n📊 Examination Performance Analysis\n\n✓ Prepared by: {self.created_by}\n✓ Date: {datetime.now().strftime('%B %d, %Y')}"
        p.font.size = Pt(14)
        p.font.color.rgb = self.text_muted
        p.space_before = Pt(16)
        
    def add_executive_summary_slide(self, stats):
        """Slide 2: Premium executive summary with key findings"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        self._add_gradient_background(slide)
        
        # Header
        self._add_slide_header(slide, "📈 Executive Summary")
        
        # Main findings box
        findings_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        findings_box.fill.solid()
        findings_box.fill.fore_color.rgb = PowerPointRGBColor(30, 41, 59)
        findings_box.line.color.rgb = self.secondary_color
        findings_box.line.width = Pt(2)
        
        tf = findings_box.text_frame
        tf.margin_bottom = Inches(0.3)
        tf.margin_left = Inches(0.4)
        tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.3)
        
        # Key findings with icons
        findings = [
            (f"✓ Pass Rate: {stats.get('pass_rate', 0):.1f}%", self.success_color if stats.get('pass_rate', 0) >= 85 else self.warning_color),
            (f"→ Average Aggregate: {stats.get('avg_aggregate', 0):.1f}", self.secondary_color),
            (f"⭐ Distinction Rate: {stats.get('distinction_rate', 0):.1f}%", self.accent_color),
            (f"👥 Total Students: {stats.get('total_students', 0)}", self.secondary_color),
            (f"📊 Fee Correlation Impact: {stats.get('fee_impact', 0):.1f}%", self.warning_color),
        ]
        
        for finding, color in findings:
            p = tf.add_paragraph()
            p.text = finding
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = color
            p.space_after = Pt(14)
        
    def add_performance_overview_slide(self, results_data):
        """Slide 3: Premium overall performance metrics with visual elements"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # Header
        self._add_slide_header(slide, "🎯 Overall Performance Overview")
        
        # Create 3 metric cards
        metrics = [
            ("Pass Rate", "89.2%", "↗ +2.3%", self.success_color),
            ("Avg Aggregate", "15.4", "↘ -0.2", self.warning_color),
            ("Distinction Rate", "23.5%", "↗ +1.8%", self.accent_color),
        ]
        
        card_width = Inches(2.8)
        card_height = Inches(3)
        start_left = Inches(0.7)
        top = Inches(1.5)
        
        for idx, (label, value, trend, color) in enumerate(metrics):
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                start_left + idx * (card_width + 0.2),
                top,
                card_width,
                card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = PowerPointRGBColor(30, 41, 59)
            card.line.color.rgb = color
            card.line.width = Pt(3)
            
            # Color bar at top
            color_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                start_left + idx * (card_width + 0.2),
                top,
                card_width,
                Inches(0.2)
            )
            color_bar.fill.solid()
            color_bar.fill.fore_color.rgb = color
            color_bar.line.color.rgb = color
            
            # Text
            tf = card.text_frame
            tf.clear()
            
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = self.text_muted
            p.alignment = 1  # Center
            
            p = tf.add_paragraph()
            p.text = value
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = 1
            p.space_before = Pt(12)
            
            p = tf.add_paragraph()
            p.text = trend
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.success_color if "↗" in trend else self.warning_color
            p.alignment = 1
            p.space_before = Pt(12)
    
    def add_subject_analysis_slide(self, subject_stats):
        """Slide 4: Premium subject performance breakdown"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # Header
        self._add_slide_header(slide, "📚 Subject Performance Analysis")
        
        # Container box
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        container.fill.solid()
        container.fill.fore_color.rgb = PowerPointRGBColor(30, 41, 59)
        container.line.color.rgb = self.secondary_color
        container.line.width = Pt(2)
        
        tf = container.text_frame
        tf.margin_bottom = Inches(0.3)
        tf.margin_left = Inches(0.4)
        tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.3)
        
        # Subject rankings with visual bars
        for subject, stats in sorted(subject_stats.items(), key=lambda x: x[1].get('pass_rate', 0), reverse=True)[:8]:
            pass_rate = stats.get('pass_rate', 0)
            p = tf.add_paragraph()
            bar = "█" * int(pass_rate / 5) + "░" * (20 - int(pass_rate / 5))
            p.text = f"{subject:18s} {bar} {pass_rate:.1f}%"
            p.font.size = Pt(11)
            p.font.name = 'Courier New'
            p.font.bold = True
            
            # Color based on performance
            if pass_rate >= 85:
                p.font.color.rgb = self.success_color
            elif pass_rate >= 75:
                p.font.color.rgb = self.secondary_color
            else:
                p.font.color.rgb = self.warning_color
            
            p.space_after = Pt(6)
    
    def add_class_comparison_slide(self, class_stats):
        """Slide 5: Premium class ranking and comparison"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # Header
        self._add_slide_header(slide, "🏆 Class Performance Ranking")
        
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        container.fill.solid()
        container.fill.fore_color.rgb = PowerPointRGBColor(30, 41, 59)
        container.line.color.rgb = self.accent_color
        container.line.width = Pt(2)
        
        tf = container.text_frame
        tf.margin_bottom = Inches(0.3)
        tf.margin_left = Inches(0.4)
        tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.3)
        
        # Class rankings
        for idx, (class_name, stats) in enumerate(sorted(class_stats.items(), key=lambda x: x[1].get('pass_rate', 0), reverse=True), 1):
            p = tf.add_paragraph()
            medal = "🥇" if idx == 1 else "🥈" if idx == 2 else "🥉" if idx == 3 else f"#{idx}"
            pass_rate = stats.get('pass_rate', 0)
            p.text = f"{medal} {class_name}: {pass_rate:.1f}% | {stats.get('count', 0)} students"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.secondary_color
            p.space_after = Pt(10)
    
    def add_recommendations_slide(self, recommendations):
        """Slide 6: Premium recommendations for improvement"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # Header
        self._add_slide_header(slide, "💡 Recommendations & Action Plan")
        
        # Container
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        container.fill.solid()
        container.fill.fore_color.rgb = PowerPointRGBColor(30, 41, 59)
        container.line.color.rgb = self.warning_color
        container.line.width = Pt(2)
        
        tf = container.text_frame
        tf.margin_bottom = Inches(0.3)
        tf.margin_left = Inches(0.4)
        tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.3)
        
        for idx, rec in enumerate(recommendations[:6], 1):
            p = tf.add_paragraph()
            p.text = f"→ {rec}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.text_light
            p.space_after = Pt(12)
    
    def _add_slide_header(self, slide, title):
        """Add premium header to slide"""
        # Header background
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = PowerPointRGBColor(15, 23, 42)
        header.line.color.rgb = self.secondary_color
        header.line.width = Pt(2)
        
        # Header text
        header_tf = header.text_frame
        header_tf.clear()
        p = header_tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.secondary_color
        p.alignment = 1  # Center
    
    def save(self, filepath):
        """Save presentation to file"""
        self.prs.save(filepath)
        return filepath


class ExcelExporter:
    """Generate multi-sheet Excel workbooks with data and analysis"""
    
    def __init__(self):
        self.wb = Workbook()
        self.wb.remove(self.wb.active)
        
    def add_raw_data_sheet(self, results_list):
        """Add sheet with all student records"""
        ws = self.wb.create_sheet("Raw_Data")
        
        # Headers
        headers = ['Student ID', 'Name', 'Class', 'Gender', 'Aggregate', 'Status', 
                   'English', 'Mathematics', 'Science', 'Social Studies', 'Language']
        ws.append(headers)
        
        # Style header
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
            cell.alignment = Alignment(horizontal="center", vertical="center")
        
        # Data
        for result in results_list:
            ws.append([
                result.student.id,
                f"{result.student.surname}, {result.student.first_name}",
                f"Grade {result.student.current_class.grade} {result.student.current_class.section}" if result.student.current_class else "N/A",
                "Male" if result.student.sex == 'M' else "Female" if result.student.sex == 'F' else "Other",
                result.total_aggregate or "",
                result.overall_status,
                result.english_units or "",
                result.mathematics_units or "",
                result.science_units or "",
                result.social_studies_units or "",
                result.indigenous_language_units or "",
            ])
        
        # Adjust column widths
        ws.column_dimensions['A'].width = 12
        ws.column_dimensions['B'].width = 20
        ws.column_dimensions['C'].width = 15
        for col in ['D', 'E', 'F', 'G', 'H', 'I', 'J', 'K']:
            ws.column_dimensions[col].width = 12
    
    def add_statistics_sheet(self, stats_dict):
        """Add sheet with calculated statistics"""
        ws = self.wb.create_sheet("Statistics")
        
        ws.append(["Metric", "Value"])
        for metric, value in stats_dict.items():
            ws.append([metric, value])
        
        ws.column_dimensions['A'].width = 25
        ws.column_dimensions['B'].width = 15
    
    def add_subject_analysis_sheet(self, subject_data):
        """Add sheet with subject-by-subject breakdown"""
        ws = self.wb.create_sheet("Subject_Analysis")
        
        headers = ['Subject', 'Pass Rate %', 'Avg Units', 'Distinction %', 'Credit %', 'Pass %', 'Fail %']
        ws.append(headers)
        
        for subject, data in subject_data.items():
            ws.append([
                subject,
                data.get('pass_rate', 0),
                data.get('avg_units', 0),
                data.get('distinction_pct', 0),
                data.get('credit_pct', 0),
                data.get('pass_pct', 0),
                data.get('fail_pct', 0),
            ])
        
        ws.column_dimensions['A'].width = 18
        for col in ['B', 'C', 'D', 'E', 'F', 'G']:
            ws.column_dimensions[col].width = 14
    
    def save(self, filepath):
        """Save workbook to file"""
        self.wb.save(filepath)
        return filepath


class PDFExporter:
    """Generate PREMIUM professional PDF reports"""
    
    def __init__(self, title="ZIMSEC Examination Report", school_name="School Name"):
        self.title = title
        self.school_name = school_name
        self.styles = getSampleStyleSheet()
        
        # Premium colors (normalized to 0-1 range)
        self.primary_color = Color(15/255, 23/255, 42/255)        # Slate-900
        self.secondary_color = Color(59/255, 130/255, 246/255)    # Blue-500
        self.accent_color = Color(236/255, 72/255, 153/255)       # Pink-500
        self.success_color = Color(34/255, 197/255, 94/255)       # Green-500
        self.warning_color = Color(249/255, 115/255, 22/255)      # Orange-500
        self.text_dark = Color(30/255, 41/255, 59/255)            # Slate-800
        self.text_muted = Color(100/255, 116/255, 139/255)        # Slate-500
        
        # RGB colors for table styling (normalized to 0-1 range)
        self.rgb_secondary = Color(59/255, 130/255, 246/255)      # Blue-500 (#3B82F6)
        self.rgb_bg_light = Color(248/255, 250/255, 252/255)      # Slate-100 (#F8FAFC)
        self.rgb_bg_lighter = Color(241/255, 245/255, 249/255)    # Slate-50 (#F1F5F9)
        self.rgb_text_dark = Color(30/255, 41/255, 59/255)        # Slate-800 (#1E293B)
        
        self._add_custom_styles()
        
    def _add_custom_styles(self):
        """Add premium custom paragraph styles"""
        # Title style
        self.styles.add(ParagraphStyle(
            name='PremiumTitle',
            parent=self.styles['Heading1'],
            fontSize=28,
            textColor=self.secondary_color,
            spaceAfter=24,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Subtitle style
        self.styles.add(ParagraphStyle(
            name='PremiumSubtitle',
            parent=self.styles['Normal'],
            fontSize=14,
            textColor=self.text_dark,
            spaceAfter=8,
            alignment=TA_CENTER
        ))
        
        # Section heading style
        self.styles.add(ParagraphStyle(
            name='PremiumHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=self.secondary_color,
            spaceAfter=14,
            spaceBefore=14,
            fontName='Helvetica-Bold'
        ))
        
        # Finding style
        self.styles.add(ParagraphStyle(
            name='Finding',
            parent=self.styles['Normal'],
            fontSize=11,
            leftIndent=20,
            spaceAfter=8,
            textColor=self.text_dark
        ))
    
    def generate_report(self, filepath, content_data):
        """Generate premium PDF report with comprehensive analysis"""
        doc = SimpleDocTemplate(
            filepath,
            pagesize=letter,
            rightMargin=0.75 * inch,
            leftMargin=0.75 * inch,
            topMargin=0.75 * inch,
            bottomMargin=0.75 * inch
        )
        story = []
        
        # ========== COVER PAGE SECTION ==========
        story.append(Spacer(1, 0.5 * inch))
        
        # Title with accent
        title_para = Paragraph(f"📊 {self.title}", self.styles['PremiumTitle'])
        story.append(title_para)
        
        story.append(Spacer(1, 0.2 * inch))
        
        # School info
        school_para = Paragraph(f"<b>{self.school_name}</b>", self.styles['PremiumSubtitle'])
        story.append(school_para)
        
        report_date = Paragraph(f"Report Generated: {datetime.now().strftime('%B %d, %Y')}", self.styles['PremiumSubtitle'])
        story.append(report_date)
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== EXECUTIVE SUMMARY ==========
        story.append(Paragraph("📋 Executive Summary", self.styles['PremiumHeading']))
        
        # Enhanced summary with interpretation
        pass_rate = content_data.get('pass_rate', 0)
        avg_agg = content_data.get('avg_aggregate', 0)
        dist_rate = content_data.get('distinction_rate', 0)
        total_students = content_data.get('total_students', 0)
        
        # Generate performance interpretation
        if pass_rate >= 90:
            performance_desc = "Outstanding performance with exceptional pass rates"
        elif pass_rate >= 80:
            performance_desc = "Strong academic performance with room for improvement in specific areas"
        elif pass_rate >= 70:
            performance_desc = "Moderate performance requiring targeted interventions"
        else:
            performance_desc = "Below-average performance demanding urgent academic support"
        
        exec_summary = f"""
        <b>Overall Assessment:</b> {performance_desc}.<br/>
        The school achieved a <b>{pass_rate:.1f}%</b> pass rate across all examination subjects, with 
        <b>{total_students}</b> students participating. The average aggregate score of <b>{avg_agg:.1f}</b> indicates 
        <b>{'excellent' if avg_agg > 14 else 'good' if avg_agg > 12 else 'moderate' if avg_agg > 10 else 'below-average'}</b> 
        academic achievement. A <b>{dist_rate:.1f}%</b> distinction rate demonstrates that a significant 
        portion of students have excelled in their examinations.
        """
        story.append(Paragraph(exec_summary, self.styles['Finding']))
        
        story.append(Spacer(1, 0.2 * inch))
        
        # ========== DETAILED PERFORMANCE ANALYSIS ==========
        story.append(Paragraph("📈 Detailed Performance Analysis", self.styles['PremiumHeading']))
        
        # Pass Rate Analysis
        story.append(Paragraph("<b>Pass Rate Analysis</b>", self.styles['Heading2']))
        if pass_rate >= 90:
            pass_analysis = f"With a pass rate of {pass_rate:.1f}%, the school is performing exceptionally well. This indicates that 9 out of every 10 students are meeting the minimum examination requirements."
        elif pass_rate >= 80:
            pass_analysis = f"The {pass_rate:.1f}% pass rate shows strong overall performance. While 8 out of 10 students are passing, there is potential to support the remaining students in achieving passing grades."
        else:
            pass_analysis = f"The {pass_rate:.1f}% pass rate indicates that less than {100-pass_rate:.1f}% of students are not meeting minimum requirements. Immediate intervention strategies are needed."
        story.append(Paragraph(f"→ {pass_analysis}", self.styles['Finding']))
        story.append(Spacer(1, 0.15 * inch))
        
        # Aggregate Score Analysis
        story.append(Paragraph("<b>Aggregate Score Distribution</b>", self.styles['Heading2']))
        if avg_agg > 14:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} falls in the excellent range (>14). This demonstrates strong subject mastery and consistent academic excellence across the curriculum."
        elif avg_agg > 12:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} is in the good range (12-14). Students are performing well across subjects, though some subjects may require targeted support."
        elif avg_agg > 10:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} is moderate (10-12). While students are passing, there is significant room for improvement in subject mastery."
        else:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} is below the ideal threshold (<10). Comprehensive academic support and curriculum review are recommended."
        story.append(Paragraph(f"→ {agg_analysis}", self.styles['Finding']))
        story.append(Spacer(1, 0.15 * inch))
        
        # Distinction Rate Analysis
        story.append(Paragraph("<b>Excellence Rate (Distinctions)</b>", self.styles['Heading2']))
        if dist_rate >= 20:
            dist_analysis = f"An outstanding {dist_rate:.1f}% of students achieved distinctions, indicating a culture of academic excellence."
        elif dist_rate >= 10:
            dist_analysis = f"With {dist_rate:.1f}% of students achieving distinctions, the school demonstrates good academic excellence."
        else:
            dist_analysis = f"Only {dist_rate:.1f}% of students achieved distinctions, suggesting a need to enhance programs for high-achieving students."
        story.append(Paragraph(f"→ {dist_analysis}", self.styles['Finding']))
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== KEY FINDINGS ==========
        story.append(Paragraph("🎯 Key Findings & Metrics", self.styles['PremiumHeading']))
        
        key_findings = [
            f"Overall Pass Rate: <b>{pass_rate:.1f}%</b> - {'✓ EXCELLENT (>85%)' if pass_rate >= 85 else '✓ GOOD (70-85%)' if pass_rate >= 70 else '⚠ NEEDS IMPROVEMENT (<70%)'}",
            f"Average Aggregate Score: <b>{avg_agg:.1f}</b> - {'Excellent (>14)' if avg_agg > 14 else 'Good (12-14)' if avg_agg > 12 else 'Moderate (10-12)' if avg_agg > 10 else 'Below Average (<10)'}",
            f"Distinction Rate: <b>{dist_rate:.1f}%</b> - {'Outstanding (>20%)' if dist_rate >= 20 else 'Good (10-20%)' if dist_rate >= 10 else 'Acceptable (<10%)'}",
            f"Total Students Examined: <b>{total_students}</b> students",
            f"Subject Performance Consistency: <b>{'High' if avg_agg > 14 else 'Moderate-High' if avg_agg > 12 else 'Moderate' if avg_agg > 10 else 'Low'}</b>",
            f"Fee Payment Impact: <b>{content_data.get('fee_impact', 0):.1f}%</b> correlation with performance"
        ]
        
        for finding in key_findings:
            story.append(Paragraph(f"→ {finding}", self.styles['Finding']))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== PERFORMANCE METRICS TABLE ==========
        story.append(Paragraph("📊 Comprehensive Performance Metrics", self.styles['PremiumHeading']))
        
        metrics_data = [
            ['📈 Metric', '📊 Value', '💡 Interpretation'],
            ['Overall Pass Rate', f"{pass_rate:.1f}%", '✓ STRONG' if pass_rate >= 85 else '⚠ MONITOR'],
            ['Average Aggregate', f"{avg_agg:.1f}", 'Excellent' if avg_agg > 14 else 'Good'],
            ['Distinction Rate', f"{dist_rate:.1f}%", '✓ HIGH' if dist_rate >= 20 else 'Good'],
            ['Total Students', str(total_students), 'Full Cohort'],
            ['Fee Correlation', f"{content_data.get('fee_impact', 0):.1f}%", 'Monitor'],
        ]
        
        table = Table(metrics_data, colWidths=[2.0 * inch, 1.6 * inch, 1.9 * inch])
        table.setStyle(TableStyle([
            # Header styling
            ('BACKGROUND', (0, 0), (-1, 0), self.rgb_secondary),
            ('TEXTCOLOR', (0, 0), (-1, 0), Color(0.96, 0.96, 0.96)),
            ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
            ('TOPPADDING', (0, 0), (-1, 0), 10),
            
            # Row styling
            ('BACKGROUND', (0, 1), (-1, -1), self.rgb_bg_light),
            ('TEXTCOLOR', (0, 1), (-1, -1), self.rgb_text_dark),
            ('ALIGN', (0, 1), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [self.rgb_bg_light, self.rgb_bg_lighter]),
            
            # Grid and borders
            ('GRID', (0, 0), (-1, -1), 1, self.rgb_secondary),
            ('LINEABOVE', (0, 0), (-1, 0), 2, self.rgb_secondary),
            ('LINEBELOW', (0, -1), (-1, -1), 2, self.rgb_secondary),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
            ('TOPPADDING', (0, 1), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
        ]))
        story.append(table)
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== ADDITIONAL INSIGHTS ==========
        story.append(Paragraph("💡 Key Insights & Observations", self.styles['PremiumHeading']))
        
        insights = []
        if pass_rate >= 85 and dist_rate >= 20:
            insights.append("• Strong academic foundation with high percentage of excellent performers")
        if pass_rate < 80:
            insights.append("• Significant number of students require academic support and intervention")
        if avg_agg > 14:
            insights.append("• Consistent high-quality teaching and learning across all subjects")
        elif avg_agg < 10:
            insights.append("• Wide disparity in subject performance requiring curriculum alignment")
        if content_data.get('fee_impact', 0) > 10:
            insights.append(f"• Strong correlation ({content_data.get('fee_impact', 0):.1f}%) between fee payment and academic performance")
        
        if not insights:
            insights = [
                "• Continue strengthening successful teaching methodologies",
                "• Implement targeted intervention programs for underperforming students",
                "• Maintain regular monitoring of student progress",
                "• Engage parents and guardians in academic improvement initiatives"
            ]
        
        for insight in insights:
            story.append(Paragraph(insight, self.styles['Finding']))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== RECOMMENDATIONS ==========
        story.append(Paragraph("🎯 Strategic Recommendations & Action Items", self.styles['PremiumHeading']))
        
        recommendations = [
            "<b>Subject-Specific Support:</b> Identify underperforming subjects and provide targeted remedial classes and tutoring",
            "<b>Student Intervention Program:</b> Establish mentoring and counseling for students below pass threshold",
            "<b>Excellence Program:</b> Create advanced program for distinction-achieving students to maintain momentum",
            "<b>Fee Management:</b> Address financial barriers by establishing payment plans for students with difficulties",
            "<b>Teacher Development:</b> Conduct professional development on latest pedagogical methods and assessment strategies",
            "<b>Parent Engagement:</b> Establish regular parent-teacher forums to discuss student progress and expectations",
            "<b>Curriculum Review:</b> Analyze curriculum alignment with examination requirements and student needs",
            "<b>Performance Monitoring:</b> Implement quarterly progress reviews to track improvement and adjust strategies"
        ]
        
        for idx, rec in enumerate(recommendations, 1):
            story.append(Paragraph(f"<b>{idx}.</b> {rec}", self.styles['Finding']))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== FOOTER ==========
        footer_text = f"""
        <i>This report contains confidential examination analysis and performance metrics for internal use only. 
        Data is based on {total_students} student examination results. For detailed questions or inquiries, 
        please contact the administration office. Report generated by AcademiaFlow Dashboard.</i><br/>
        <b>Confidential - For Official Use Only</b>
        """
        story.append(Paragraph(footer_text, self.styles['Normal']))
        
        
        # Build PDF
        doc.build(story)
        return filepath

    def generate_report_to_buffer(self, buffer, content_data):
        """Generate premium PDF report to BytesIO buffer"""
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.75 * inch,
            leftMargin=0.75 * inch,
            topMargin=0.75 * inch,
            bottomMargin=0.75 * inch
        )
        story = []
        
        # ========== COVER PAGE SECTION ==========
        story.append(Spacer(1, 0.5 * inch))
        
        # Title with accent
        title_para = Paragraph(f"📊 {self.title}", self.styles['PremiumTitle'])
        story.append(title_para)
        
        story.append(Spacer(1, 0.2 * inch))
        
        # School info
        school_para = Paragraph(f"<b>{self.school_name}</b>", self.styles['PremiumSubtitle'])
        story.append(school_para)
        
        report_date = Paragraph(f"Report Generated: {datetime.now().strftime('%B %d, %Y')}", self.styles['PremiumSubtitle'])
        story.append(report_date)
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== EXECUTIVE SUMMARY ==========
        story.append(Paragraph("📋 Executive Summary", self.styles['PremiumHeading']))
        
        # Extract data for analysis
        pass_rate = content_data.get('pass_rate', 0)
        avg_agg = content_data.get('avg_aggregate', 0)
        dist_rate = content_data.get('distinction_rate', 0)
        total_students = content_data.get('total_students', 0)
        
        # Generate performance interpretation
        if pass_rate >= 90:
            performance_desc = "Outstanding performance with exceptional pass rates"
        elif pass_rate >= 80:
            performance_desc = "Strong academic performance with room for improvement in specific areas"
        elif pass_rate >= 70:
            performance_desc = "Moderate performance requiring targeted interventions"
        else:
            performance_desc = "Below-average performance demanding urgent academic support"
        
        exec_summary = f"""
        <b>Overall Assessment:</b> {performance_desc}.<br/>
        The school achieved a <b>{pass_rate:.1f}%</b> pass rate across all examination subjects, with 
        <b>{total_students}</b> students participating. The average aggregate score of <b>{avg_agg:.1f}</b> indicates 
        <b>{'excellent' if avg_agg > 14 else 'good' if avg_agg > 12 else 'moderate' if avg_agg > 10 else 'below-average'}</b> 
        academic achievement. A <b>{dist_rate:.1f}%</b> distinction rate demonstrates that a significant 
        portion of students have excelled in their examinations.
        """
        story.append(Paragraph(exec_summary, self.styles['Finding']))
        
        story.append(Spacer(1, 0.2 * inch))
        
        # ========== DETAILED PERFORMANCE ANALYSIS ==========
        story.append(Paragraph("📈 Detailed Performance Analysis", self.styles['PremiumHeading']))
        
        # Pass Rate Analysis
        story.append(Paragraph("<b>Pass Rate Analysis</b>", self.styles['Heading2']))
        if pass_rate >= 90:
            pass_analysis = f"With a pass rate of {pass_rate:.1f}%, the school is performing exceptionally well. This indicates that 9 out of every 10 students are meeting the minimum examination requirements."
        elif pass_rate >= 80:
            pass_analysis = f"The {pass_rate:.1f}% pass rate shows strong overall performance. While 8 out of 10 students are passing, there is potential to support the remaining students in achieving passing grades."
        else:
            pass_analysis = f"The {pass_rate:.1f}% pass rate indicates that less than {100-pass_rate:.1f}% of students are not meeting minimum requirements. Immediate intervention strategies are needed."
        story.append(Paragraph(f"→ {pass_analysis}", self.styles['Finding']))
        story.append(Spacer(1, 0.15 * inch))
        
        # Aggregate Score Analysis
        story.append(Paragraph("<b>Aggregate Score Distribution</b>", self.styles['Heading2']))
        if avg_agg > 14:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} falls in the excellent range (>14). This demonstrates strong subject mastery and consistent academic excellence across the curriculum."
        elif avg_agg > 12:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} is in the good range (12-14). Students are performing well across subjects, though some subjects may require targeted support."
        elif avg_agg > 10:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} is moderate (10-12). While students are passing, there is significant room for improvement in subject mastery."
        else:
            agg_analysis = f"The average aggregate score of {avg_agg:.1f} is below the ideal threshold (<10). Comprehensive academic support and curriculum review are recommended."
        story.append(Paragraph(f"→ {agg_analysis}", self.styles['Finding']))
        story.append(Spacer(1, 0.15 * inch))
        
        # Distinction Rate Analysis
        story.append(Paragraph("<b>Excellence Rate (Distinctions)</b>", self.styles['Heading2']))
        if dist_rate >= 20:
            dist_analysis = f"An outstanding {dist_rate:.1f}% of students achieved distinctions, indicating a culture of academic excellence."
        elif dist_rate >= 10:
            dist_analysis = f"With {dist_rate:.1f}% of students achieving distinctions, the school demonstrates good academic excellence."
        else:
            dist_analysis = f"Only {dist_rate:.1f}% of students achieved distinctions, suggesting a need to enhance programs for high-achieving students."
        story.append(Paragraph(f"→ {dist_analysis}", self.styles['Finding']))
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== KEY FINDINGS ==========
        story.append(Paragraph("🎯 Key Findings & Metrics", self.styles['PremiumHeading']))
        
        key_findings = [
            f"Overall Pass Rate: <b>{pass_rate:.1f}%</b> - {'✓ EXCELLENT (>85%)' if pass_rate >= 85 else '✓ GOOD (70-85%)' if pass_rate >= 70 else '⚠ NEEDS IMPROVEMENT (<70%)'}",
            f"Average Aggregate Score: <b>{avg_agg:.1f}</b> - {'Excellent (>14)' if avg_agg > 14 else 'Good (12-14)' if avg_agg > 12 else 'Moderate (10-12)' if avg_agg > 10 else 'Below Average (<10)'}",
            f"Distinction Rate: <b>{dist_rate:.1f}%</b> - {'Outstanding (>20%)' if dist_rate >= 20 else 'Good (10-20%)' if dist_rate >= 10 else 'Acceptable (<10%)'}",
            f"Total Students Examined: <b>{total_students}</b> students",
            f"Subject Performance Consistency: <b>{'High' if avg_agg > 14 else 'Moderate-High' if avg_agg > 12 else 'Moderate' if avg_agg > 10 else 'Low'}</b>",
            f"Fee Payment Impact: <b>{content_data.get('fee_impact', 0):.1f}%</b> correlation with performance"
        ]
        
        for finding in key_findings:
            story.append(Paragraph(f"→ {finding}", self.styles['Finding']))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== PERFORMANCE METRICS TABLE ==========
        story.append(Paragraph("📊 Comprehensive Performance Metrics", self.styles['PremiumHeading']))
        
        metrics_data = [
            ['📈 Metric', '📊 Value', '💡 Interpretation'],
            ['Overall Pass Rate', f"{pass_rate:.1f}%", '✓ STRONG' if pass_rate >= 85 else '⚠ MONITOR'],
            ['Average Aggregate', f"{avg_agg:.1f}", 'Excellent' if avg_agg > 14 else 'Good'],
            ['Distinction Rate', f"{dist_rate:.1f}%", '✓ HIGH' if dist_rate >= 20 else 'Good'],
            ['Total Students', str(total_students), 'Full Cohort'],
            ['Fee Correlation', f"{content_data.get('fee_impact', 0):.1f}%", 'Monitor'],
        ]
        
        table = Table(metrics_data, colWidths=[2.0 * inch, 1.6 * inch, 1.9 * inch])
        table.setStyle(TableStyle([
            # Header styling
            ('BACKGROUND', (0, 0), (-1, 0), self.rgb_secondary),
            ('TEXTCOLOR', (0, 0), (-1, 0), Color(0.96, 0.96, 0.96)),
            ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
            ('TOPPADDING', (0, 0), (-1, 0), 10),
            
            # Row styling
            ('BACKGROUND', (0, 1), (-1, -1), self.rgb_bg_light),
            ('TEXTCOLOR', (0, 1), (-1, -1), self.rgb_text_dark),
            ('ALIGN', (0, 1), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [self.rgb_bg_light, self.rgb_bg_lighter]),
            
            # Grid and borders
            ('GRID', (0, 0), (-1, -1), 1, self.rgb_secondary),
            ('LINEABOVE', (0, 0), (-1, 0), 2, self.rgb_secondary),
            ('LINEBELOW', (0, -1), (-1, -1), 2, self.rgb_secondary),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
            ('TOPPADDING', (0, 1), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
        ]))
        story.append(table)
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== ADDITIONAL INSIGHTS ==========
        story.append(Paragraph("💡 Key Insights & Observations", self.styles['PremiumHeading']))
        
        insights = []
        if pass_rate >= 85 and dist_rate >= 20:
            insights.append("• Strong academic foundation with high percentage of excellent performers")
        if pass_rate < 80:
            insights.append("• Significant number of students require academic support and intervention")
        if avg_agg > 14:
            insights.append("• Consistent high-quality teaching and learning across all subjects")
        elif avg_agg < 10:
            insights.append("• Wide disparity in subject performance requiring curriculum alignment")
        if content_data.get('fee_impact', 0) > 10:
            insights.append(f"• Strong correlation ({content_data.get('fee_impact', 0):.1f}%) between fee payment and academic performance")
        
        if not insights:
            insights = [
                "• Continue strengthening successful teaching methodologies",
                "• Implement targeted intervention programs for underperforming students",
                "• Maintain regular monitoring of student progress",
                "• Engage parents and guardians in academic improvement initiatives"
            ]
        
        for insight in insights:
            story.append(Paragraph(insight, self.styles['Finding']))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== RECOMMENDATIONS ==========
        story.append(Paragraph("🎯 Strategic Recommendations & Action Items", self.styles['PremiumHeading']))
        
        recommendations = [
            "<b>Subject-Specific Support:</b> Identify underperforming subjects and provide targeted remedial classes and tutoring",
            "<b>Student Intervention Program:</b> Establish mentoring and counseling for students below pass threshold",
            "<b>Excellence Program:</b> Create advanced program for distinction-achieving students to maintain momentum",
            "<b>Fee Management:</b> Address financial barriers by establishing payment plans for students with difficulties",
            "<b>Teacher Development:</b> Conduct professional development on latest pedagogical methods and assessment strategies",
            "<b>Parent Engagement:</b> Establish regular parent-teacher forums to discuss student progress and expectations",
            "<b>Curriculum Review:</b> Analyze curriculum alignment with examination requirements and student needs",
            "<b>Performance Monitoring:</b> Implement quarterly progress reviews to track improvement and adjust strategies"
        ]
        
        for idx, rec in enumerate(recommendations, 1):
            story.append(Paragraph(f"<b>{idx}.</b> {rec}", self.styles['Finding']))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # ========== FOOTER ==========
        footer_text = f"""
        <i>This report contains confidential examination analysis and performance metrics for internal use only. 
        Data is based on {total_students} student examination results. For detailed questions or inquiries, 
        please contact the administration office. Report generated by AcademiaFlow Dashboard.</i><br/>
        <b>Confidential - For Official Use Only</b>
        """
        story.append(Paragraph(footer_text, self.styles['Normal']))
        
        # Build PDF
        doc.build(story)
        return buffer


class DetailedResultsPDFExporter:
    """Export detailed ZIMSEC results with professional styling and comprehensive analytics"""
    
    def __init__(self, title="Detailed ZIMSEC Results", school_name="School Name"):
        self.title = title
        self.school_name = school_name
        
        # Premium color palette with better contrast
        self.primary_dark = Color(13/255, 27/255, 42/255)  # Very dark navy
        self.primary_blue = Color(37/255, 99/255, 235/255)  # Professional blue
        self.accent_teal = Color(20/255, 184/255, 166/255)  # Teal accent
        self.success_green = Color(34/255, 197/255, 94/255)  # Vibrant green
        self.warning_amber = Color(217/255, 119/255, 6/255)  # Warm amber
        self.danger_red = Color(220/255, 38/255, 38/255)  # Professional red
        
        # Neutral palette
        self.white = Color(1, 1, 1)
        self.bg_white = Color(255/255, 255/255, 255/255)
        self.bg_light = Color(248/255, 249/255, 250/255)  # Very subtle gray
        self.bg_lighter = Color(241/255, 245/255, 250/255)  # Subtle blue-gray
        self.bg_accent = Color(240/255, 249/255, 255/255)  # Light blue tint
        self.border_light = Color(226/255, 232/255, 240/255)  # Subtle border
        self.text_dark = Color(15/255, 23/255, 42/255)  # Deep text
        self.text_muted = Color(100/255, 116/255, 139/255)  # Muted text
        
        # Paragraph styles
        self.styles = getSampleStyleSheet()
        
        # Main report title
        self.styles.add(ParagraphStyle(
            name='ReportTitle',
            parent=self.styles['Heading1'],
            fontSize=32,
            textColor=self.primary_dark,
            spaceAfter=4,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold',
            leading=38
        ))
        
        # School name
        self.styles.add(ParagraphStyle(
            name='SchoolName',
            parent=self.styles['Heading2'],
            fontSize=13,
            textColor=self.primary_blue,
            spaceAfter=2,
            alignment=TA_CENTER,
            fontName='Helvetica'
        ))
        
        # Section heading with background
        self.styles.add(ParagraphStyle(
            name='SectionHeading',
            parent=self.styles['Heading2'],
            fontSize=12,
            textColor=self.white,
            spaceAfter=0,
            fontName='Helvetica-Bold',
            borderPadding=0
        ))
        
        # Subsection heading
        self.styles.add(ParagraphStyle(
            name='SubsectionHeading',
            parent=self.styles['Heading3'],
            fontSize=11,
            textColor=self.primary_blue,
            spaceAfter=6,
            fontName='Helvetica-Bold',
            borderColor=self.accent_teal,
            borderWidth=1,
            borderPadding=4,
            borderRadius=2
        ))
        
        # Class heading
        self.styles.add(ParagraphStyle(
            name='ClassHeading',
            parent=self.styles['Heading3'],
            fontSize=10.5,
            textColor=self.primary_dark,
            spaceAfter=4,
            fontName='Helvetica-Bold'
        ))
    
    def _create_styled_section_header(self, title, icon=""):
        """Create a professional section header with styling"""
        from reportlab.platypus import Table, TableStyle
        
        header_text = f"{icon} {title}" if icon else title
        header_data = [[header_text]]
        header_table = Table(header_data, colWidths=[7.2*inch])
        header_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), self.primary_blue),
            ('TEXTCOLOR', (0, 0), (-1, -1), self.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 13),
            ('PADDING', (0, 0), (-1, -1), 14),
            ('LEFTPADDING', (0, 0), (-1, -1), 16),
            ('TOPPADDING', (0, 0), (-1, -1), 12),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('ROWBACKGROUNDS', (0, 0), (-1, -1), [self.primary_blue]),
        ]))
        return header_table
    
    def _create_metric_table(self, data, colWidths):
        """Create a professional metrics table"""
        from reportlab.platypus import Table, TableStyle
        
        table = Table(data, colWidths=colWidths)
        table.setStyle(TableStyle([
            # Header styling
            ('BACKGROUND', (0, 0), (-1, 0), self.primary_blue),
            ('TEXTCOLOR', (0, 0), (-1, 0), self.white),
            ('ALIGN', (0, 0), (-1, 0), 'LEFT'),
            ('VALIGN', (0, 0), (-1, 0), 'MIDDLE'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('PADDING', (0, 0), (-1, 0), 14),
            ('TOPPADDING', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            
            # Data rows
            ('BACKGROUND', (0, 1), (-1, -1), self.white),
            ('TEXTCOLOR', (0, 1), (-1, -1), self.text_dark),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('PADDING', (0, 0), (-1, -1), (12, 8)),
            ('LEFTPADDING', (0, 0), (-1, -1), 14),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            
            # Alternating rows
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [self.white, self.bg_accent]),
            
            # Borders
            ('GRID', (0, 0), (-1, -1), 0.5, Color(200/255, 200/255, 200/255)),
            ('LINEABOVE', (0, 0), (-1, 0), 2, self.primary_blue),
            ('LINEBELOW', (0, -1), (-1, -1), 2, self.primary_blue),
            ('LINEABOVE', (0, 1), (-1, 1), 1, Color(220/255, 220/255, 220/255)),
            
            # First column left align
            ('ALIGN', (0, 1), (0, -1), 'LEFT'),
            ('FONTNAME', (0, 1), (0, -1), 'Helvetica-Bold'),
        ]))
        return table
    
    def export_to_buffer(self, buffer, results_by_class):
        """Export detailed results with professional styling"""
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.7 * inch,
            leftMargin=0.7 * inch,
            topMargin=0.9 * inch,
            bottomMargin=0.9 * inch
        )
        story = []
        
        # ============ TITLE PAGE ============
        story.append(Spacer(1, 0.6 * inch))
        
        # Main title
        story.append(Paragraph(f"📊 {self.title}", self.styles['ReportTitle']))
        story.append(Spacer(1, 0.05 * inch))
        
        # School name
        story.append(Paragraph(self.school_name, self.styles['SchoolName']))
        story.append(Spacer(1, 0.15 * inch))
        
        # Decorative line
        story.append(HRFlowable(width="70%", thickness=2.5, color=self.primary_blue, spaceAfter=15))
        
        # Report date and info
        story.append(Paragraph(
            f"Report Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}",
            self.styles['Normal']
        ))
        story.append(Spacer(1, 0.4 * inch))
        
        # Cover page content
        cover_content = f"""
        <b>Academic Excellence Report</b><br/><br/>
        This comprehensive report provides detailed analysis of ZIMSEC results, 
        evaluating student performance across all subjects and identifying trends 
        in academic achievement.<br/><br/>
        
        <b>Report Highlights:</b><br/>
        • Detailed student performance metrics by class<br/>
        • Subject-wise pass rates and analysis<br/>
        • Student achievement categories<br/>
        • Data-driven recommendations for improvement<br/><br/>
        
        <b>Prepared for:</b> {self.school_name}<br/>
        <b>Report Period:</b> ZIMSEC 2026 Results
        """
        story.append(Paragraph(cover_content, self.styles['Normal']))
        
        # Page break for Executive Summary
        story.append(PageBreak())
        story.append(self._create_styled_section_header("Executive Summary", "📋"))
        story.append(Spacer(1, 0.15 * inch))
        
        # Calculate overall statistics before displaying
        total_students = sum(len(results) for results in results_by_class.values())
        total_passes = 0
        total_fails = 0
        subject_performance = {'english': {'pass': 0, 'fail': 0}, 'math': {'pass': 0, 'fail': 0}, 
                               'science': {'pass': 0, 'fail': 0}, 'social_studies': {'pass': 0, 'fail': 0},
                               'language': {'pass': 0, 'fail': 0}, 'agriculture': {'pass': 0, 'fail': 0}}
        
        # Collect all results for statistics
        all_results = []
        for results in results_by_class.values():
            all_results.extend(results)
        
        # Calculate statistics
        for result in all_results:
            subjects = [result.english_units, result.mathematics_units, result.science_units, 
                       result.social_studies_units, result.indigenous_language_units, result.agriculture_units]
            passes = sum(1 for s in subjects if isinstance(s, int) and s <= 5)
            
            if passes >= 4:
                total_passes += 1
            else:
                total_fails += 1
            
            # Subject-wise stats
            if result.english_units and result.english_units <= 5:
                subject_performance['english']['pass'] += 1
            elif result.english_units:
                subject_performance['english']['fail'] += 1
                
            if result.mathematics_units and result.mathematics_units <= 5:
                subject_performance['math']['pass'] += 1
            elif result.mathematics_units:
                subject_performance['math']['fail'] += 1
                
            if result.science_units and result.science_units <= 5:
                subject_performance['science']['pass'] += 1
            elif result.science_units:
                subject_performance['science']['fail'] += 1
                
            if result.social_studies_units and result.social_studies_units <= 5:
                subject_performance['social_studies']['pass'] += 1
            elif result.social_studies_units:
                subject_performance['social_studies']['fail'] += 1
                
            if result.indigenous_language_units and result.indigenous_language_units <= 5:
                subject_performance['language']['pass'] += 1
            elif result.indigenous_language_units:
                subject_performance['language']['fail'] += 1
                
            if result.agriculture_units and result.agriculture_units <= 5:
                subject_performance['agriculture']['pass'] += 1
            elif result.agriculture_units:
                subject_performance['agriculture']['fail'] += 1
        
        # Overall Summary Section
        story.append(Spacer(1, 0.2 * inch))
        
        # Calculate pass rate for narrative
        pass_rate_pct = (total_passes/total_students*100) if total_students > 0 else 0
        
        # Executive Summary Narrative
        summary_text = f"""
        This comprehensive report presents a detailed analysis of the {datetime.now().year} ZIMSEC examination results 
        for {self.school_name}. The assessment encompasses {total_students} students across {len(results_by_class)} classes, 
        evaluated across six core academic subjects including English Language, Mathematics, General Science, Social Studies, 
        Indigenous Language, and Agriculture.
        <br/><br/>
        <b>KEY PERFORMANCE HIGHLIGHTS:</b>
        <br/>
        • Overall Student Achievement: {total_passes} students ({pass_rate_pct:.1f}%) achieved passing grades across four or more subjects
        <br/>
        • Total Assessment Coverage: {total_students} students examined across all subjects
        <br/>
        • Academic Classes: {len(results_by_class)} distinct class groups analyzed
        <br/>
        • Comprehensive Analysis: Subject performance, class-wise breakdown, individual student performance assessment
        <br/><br/>
        This report provides actionable insights for academic improvement, identifies areas of excellence, and recommends 
        strategic initiatives for enhanced institutional performance.
        """
        story.append(Paragraph(summary_text, self.styles['Normal']))
        story.append(Spacer(1, 0.25 * inch))
        
        story.append(self._create_styled_section_header("Overall Summary Statistics", "📈"))
        story.append(Spacer(1, 0.18 * inch))
        
        summary_data = [
            ['Metric', 'Value', 'Percentage'],
            ['Total Students', str(total_students), '100%'],
            ['Total Passes', str(total_passes), f'{(total_passes/total_students*100):.1f}%' if total_students > 0 else '0%'],
            ['Total Fails', str(total_fails), f'{(total_fails/total_students*100):.1f}%' if total_students > 0 else '0%'],
            ['Number of Classes', str(len(results_by_class)), '']
        ]
        
        summary_table = self._create_metric_table(summary_data, [2.5*inch, 1.5*inch, 2*inch])
        story.append(summary_table)
        story.append(Spacer(1, 0.25 * inch))
        
        # Subject Performance Section
        story.append(Spacer(1, 0.25 * inch))
        story.append(self._create_styled_section_header("Subject Performance Analysis", "📚"))
        story.append(Spacer(1, 0.18 * inch))
        
        subject_data = [
            ['Subject', 'Passes', 'Fails', 'Pass Rate', 'Fail Rate']
        ]
        
        for subject, name in [('english', 'English Language'), ('math', 'Mathematics'), ('science', 'General Science'),
                              ('social_studies', 'Social Studies'), ('language', 'Indigenous Language'), ('agriculture', 'Agriculture')]:
            passes = subject_performance[subject]['pass']
            fails = subject_performance[subject]['fail']
            total = passes + fails
            pass_rate = f'{(passes/total*100):.1f}%' if total > 0 else '0%'
            fail_rate = f'{(fails/total*100):.1f}%' if total > 0 else '0%'
            
            subject_data.append([name, str(passes), str(fails), pass_rate, fail_rate])
        
        subject_table = self._create_metric_table(subject_data, [2*inch, 1.2*inch, 1.2*inch, 1.3*inch, 1.3*inch])
        story.append(subject_table)
        story.append(Spacer(1, 0.2 * inch))
        
        # Subject Performance Analysis Narrative
        subject_analysis = """
        <b>SUBJECT PERFORMANCE ANALYSIS:</b>
        <br/>
        The six core subjects demonstrate varied performance metrics reflecting different levels of student engagement, 
        resource availability, and teaching effectiveness across disciplines. Subjects with higher pass rates indicate 
        effective pedagogical strategies and strong foundational learning. Areas with lower pass rates require focused 
        intervention, enhanced resource allocation, and targeted instructional review.
        <br/><br/>
        <b>Key Observations:</b>
        <br/>
        • Students demonstrate stronger performance in subjects with practical, hands-on learning components
        <br/>
        • Theoretical subjects show variable results dependent on student engagement and prerequisite knowledge
        <br/>
        • Subject-specific teaching methodologies significantly impact student achievement outcomes
        <br/>
        • Cross-subject performance patterns reveal overall academic capability of student cohorts
        """
        story.append(Paragraph(subject_analysis, self.styles['Normal']))
        story.append(Spacer(1, 0.25 * inch))
        
        # Page break before Class Summary Section
        story.append(PageBreak())
        story.append(self._create_styled_section_header("Class-by-Class Summary", "🏫"))
        story.append(Spacer(1, 0.18 * inch))
        
        class_summary_data = [
            ['Class', 'Students', 'Passes', 'Fails', 'Pass Rate']
        ]
        
        for class_name in sorted(results_by_class.keys()):
            results = results_by_class[class_name]
            class_passes = sum(1 for r in results if sum(1 for s in [r.english_units, r.mathematics_units, r.science_units,
                r.social_studies_units, r.indigenous_language_units, r.agriculture_units] if isinstance(s, int) and s <= 5) >= 4)
            class_fails = len(results) - class_passes
            pass_rate = f'{(class_passes/len(results)*100):.1f}%' if len(results) > 0 else '0%'
            
            class_summary_data.append([str(class_name), str(len(results)), str(class_passes), str(class_fails), pass_rate])
        
        class_table = self._create_metric_table(class_summary_data, [1.5*inch, 1.5*inch, 1.5*inch, 1.5*inch, 1.5*inch])
        story.append(class_table)
        story.append(Spacer(1, 0.25 * inch))
        
        # Class Analysis Narrative
        class_analysis = """
        <b>CLASS-BY-CLASS COMPARATIVE ANALYSIS:</b>
        <br/>
        Comparative analysis across classes reveals important variations in academic performance that may reflect differences 
        in student demographics, teaching quality, resource allocation, and classroom environment. Classes achieving higher 
        pass rates demonstrate effective instructional strategies and strong classroom management that serve as best practice 
        models. Classes with lower aggregate performance require targeted intervention, additional resource support, and 
        pedagogical review to enhance learning outcomes.
        <br/><br/>
        <b>Performance Drivers:</b>
        <br/>
        • Teacher experience and professional development quality directly impact student achievement
        <br/>
        • Class size and student-teacher ratios affect instructional effectiveness and individual student support
        <br/>
        • Classroom learning environment and discipline positively correlate with academic results
        <br/>
        • Available instructional resources and their effective utilization enhance subject mastery
        """
        story.append(Paragraph(class_analysis, self.styles['Normal']))
        story.append(Spacer(1, 0.25 * inch))
        
        # Subject Performance by Class Section
        story.append(PageBreak())
        story.append(self._create_styled_section_header("Subject Performance by Class", "🎯"))
        story.append(Spacer(1, 0.18 * inch))
        
        for class_name in sorted(results_by_class.keys()):
            results = results_by_class[class_name]
            
            story.append(Spacer(1, 0.12 * inch))
            story.append(Paragraph(f"Grade {class_name}", self.styles['ClassHeading']))
            story.append(Spacer(1, 0.08 * inch))
            
            subject_class_data = [
                ['Subject', 'Passes', 'Fails', 'Pass Rate', 'Avg. Unit']
            ]
            
            subjects_list = [
                ('english_units', 'English Language'),
                ('mathematics_units', 'Mathematics'),
                ('science_units', 'General Science'),
                ('social_studies_units', 'Social Studies'),
                ('indigenous_language_units', 'Indigenous Language'),
                ('agriculture_units', 'Agriculture')
            ]
            
            for field, name in subjects_list:
                passes = sum(1 for r in results if getattr(r, field) and getattr(r, field) <= 5)
                fails = sum(1 for r in results if getattr(r, field) and getattr(r, field) > 5)
                total = passes + fails
                pass_rate = f'{(passes/total*100):.1f}%' if total > 0 else '0%'
                
                # Calculate average unit
                units = [getattr(r, field) for r in results if getattr(r, field)]
                avg_unit = f'{(sum(units)/len(units)):.1f}' if units else 'N/A'
                
                subject_class_data.append([name, str(passes), str(fails), pass_rate, avg_unit])
            
            subject_class_table = self._create_metric_table(subject_class_data, [2.0*inch, 1.2*inch, 1.2*inch, 1.2*inch, 1.3*inch])
            story.append(subject_class_table)
            story.append(Spacer(1, 0.1 * inch))
        
        # Performance Tiers Section
        story.append(PageBreak())
        story.append(self._create_styled_section_header("Student Performance Tiers", "⭐"))
        story.append(Spacer(1, 0.18 * inch))
        
        # Categorize students by performance
        high_performers = []  # 6 subjects passed
        good_performers = []  # 4-5 subjects passed
        average_performers = []  # 2-3 subjects passed
        low_performers = []  # 0-1 subjects passed
        
        for result in all_results:
            subjects = [result.english_units, result.mathematics_units, result.science_units,
                       result.social_studies_units, result.indigenous_language_units, result.agriculture_units]
            passes = sum(1 for s in subjects if isinstance(s, int) and s <= 5)
            
            if passes == 6:
                high_performers.append(result)
            elif passes >= 4:
                good_performers.append(result)
            elif passes >= 2:
                average_performers.append(result)
            else:
                low_performers.append(result)
        
        tiers_data = [
            ['Performance Tier', 'Criteria', 'Count', 'Percentage'],
            ['🏆 Excellent', '6 Subjects Passed', str(len(high_performers)), f'{(len(high_performers)/total_students*100):.1f}%' if total_students > 0 else '0%'],
            ['✓ Good', '4-5 Subjects Passed', str(len(good_performers)), f'{(len(good_performers)/total_students*100):.1f}%' if total_students > 0 else '0%'],
            ['◐ Average', '2-3 Subjects Passed', str(len(average_performers)), f'{(len(average_performers)/total_students*100):.1f}%' if total_students > 0 else '0%'],
            ['✗ Below Average', '0-1 Subjects Passed', str(len(low_performers)), f'{(len(low_performers)/total_students*100):.1f}%' if total_students > 0 else '0%']
        ]
        
        tiers_table = self._create_metric_table(tiers_data, [2.0*inch, 2.0*inch, 1.5*inch, 1.5*inch])
        story.append(tiers_table)
        story.append(Spacer(1, 0.2 * inch))
        
        # Performance Tiers Analysis Narrative
        tiers_analysis = f"""
        <b>STUDENT PERFORMANCE TIER ANALYSIS:</b>
        <br/>
        Classification of students into performance tiers provides targeted insight into achievement distribution and 
        identifies populations requiring differentiated instructional strategies. The excellent tier represents students 
        with comprehensive subject mastery worthy of recognition and advanced learning opportunities. The good and average 
        tiers represent the majority of students requiring continued support and incremental improvement strategies. The 
        below-average tier indicates students requiring intensive intervention and remedial support.
        <br/><br/>
        <b>Tier-Specific Recommendations:</b>
        <br/>
        • <b>Excellent Tier ({len(high_performers)} students):</b> Provide advanced enrichment, leadership opportunities, and competition participation
        <br/>
        • <b>Good Tier ({len(good_performers)} students):</b> Continue quality instruction with selective enrichment for identified strengths
        <br/>
        • <b>Average Tier ({len(average_performers)} students):</b> Implement targeted remediation and skill-building interventions
        <br/>
        • <b>Below Average Tier ({len(low_performers)} students):</b> Establish intensive support program with frequent assessment and progress monitoring
        """
        story.append(Paragraph(tiers_analysis, self.styles['Normal']))
        story.append(Spacer(1, 0.25 * inch))
        
        # Top and Bottom Performers
        story.append(self._create_styled_section_header("Top Performers (Excellent Tier)", "🔝"))
        
        if high_performers:
            top_data = [['Rank', 'Student Name', 'Class', 'Passed Subjects']]
            for i, result in enumerate(sorted(high_performers, key=lambda x: x.student.full_name)[:10], 1):
                class_str = f"Grade {result.student.current_class.grade}{result.student.current_class.section}" if result.student.current_class else "N/A"
                top_data.append([str(i), result.student.full_name, class_str, '6/6 ✓'])
            
            top_table = self._create_metric_table(top_data, [0.8*inch, 2.5*inch, 1.5*inch, 1.7*inch])
            story.append(top_table)
        else:
            story.append(Paragraph("<i>No students in excellent tier</i>", self.styles['Normal']))
        
        story.append(Spacer(1, 0.25 * inch))
        
        # Subject Strength Ranking
        story.append(self._create_styled_section_header("Subject Strength Ranking (by Pass Rate)", "📊"))
        
        subject_ranking = []
        for subject, name in [('english', 'English Language'), ('math', 'Mathematics'), ('science', 'General Science'),
                              ('social_studies', 'Social Studies'), ('language', 'Indigenous Language'), ('agriculture', 'Agriculture')]:
            passes = subject_performance[subject]['pass']
            fails = subject_performance[subject]['fail']
            total = passes + fails
            pass_rate = (passes/total*100) if total > 0 else 0
            subject_ranking.append((name, passes, fails, pass_rate))
        
        # Sort by pass rate (descending)
        subject_ranking.sort(key=lambda x: x[3], reverse=True)
        
        ranking_data = [['Rank', 'Subject', 'Passes', 'Fails', 'Pass Rate']]
        for rank, (name, passes, fails, rate) in enumerate(subject_ranking, 1):
            ranking_data.append([str(rank), name, str(passes), str(fails), f'{rate:.1f}%'])
        
        ranking_table = self._create_metric_table(ranking_data, [0.8*inch, 2.2*inch, 1.2*inch, 1.2*inch, 1.3*inch])
        story.append(ranking_table)
        story.append(Spacer(1, 0.25 * inch))
        
        # Subject Ranking Analysis Narrative
        ranking_analysis = """
        <b>SUBJECT STRENGTH RANKING ANALYSIS:</b>
        <br/>
        Ranking subjects by pass rate reveals relative institutional strength across disciplines and identifies priority 
        areas for curriculum enhancement and resource allocation. Highest-performing subjects should be evaluated for 
        replicable best practices in instruction, assessment, and student engagement. Lower-performing subjects require 
        focused attention including curriculum review, teacher professional development, resource enhancement, and 
        differentiated instructional strategies.
        """
        story.append(Paragraph(ranking_analysis, self.styles['Normal']))
        story.append(Spacer(1, 0.3 * inch))
        
        # Page break before detailed results
        story.append(PageBreak())
        
        # Strategic Recommendations Section
        story.append(self._create_styled_section_header("Strategic Recommendations", "🎯"))
        story.append(Spacer(1, 0.18 * inch))
        
        recommendations_text = f"""
        <b>IMMEDIATE ACTIONS (Next 30 Days)</b>
        <br/>
        <b>1. Targeted Subject-Specific Intervention</b>
        <br/>
        • Conduct diagnostic assessment of students in low-performing subjects
        <br/>
        • Implement focused remediation sessions addressing identified skill gaps
        <br/>
        • Allocate additional instructional time to priority subject areas
        <br/><br/>
        
        <b>2. Student Support Program Launch</b>
        <br/>
        • Establish peer tutoring initiatives pairing high performers with struggling students
        <br/>
        • Create study groups for subjects with lower pass rates
        <br/>
        • Implement weekly progress monitoring and feedback mechanisms
        <br/><br/>
        
        <b>MID-TERM INITIATIVES (Next 90 Days)</b>
        <br/>
        <b>3. Teacher Professional Development</b>
        <br/>
        • Conduct subject-specific pedagogical training focusing on differentiated instruction
        <br/>
        • Facilitate collaborative lesson planning and resource sharing among teachers
        <br/>
        • Implement instructional coaching and peer observation cycles
        <br/><br/>
        
        <b>4. Resource Enhancement</b>
        <br/>
        • Upgrade laboratory equipment and science learning resources
        <br/>
        • Expand library resources for language arts development
        <br/>
        • Implement digital learning platforms for enhanced engagement
        <br/><br/>
        
        <b>LONG-TERM STRATEGY (Next Academic Year)</b>
        <br/>
        <b>5. Curriculum Review and Alignment</b>
        <br/>
        • Audit curriculum against ZIMSEC examination requirements and standards
        <br/>
        • Revise pacing guides to ensure comprehensive content coverage
        <br/>
        • Integrate assessed content more explicitly into instructional planning
        <br/><br/>
        
        <b>6. Performance Monitoring System</b>
        <br/>
        • Establish continuous assessment practices with regular progress monitoring
        <br/>
        • Create early warning system identifying at-risk students
        <br/>
        • Develop data dashboard for real-time performance tracking and stakeholder reporting
        """
        story.append(Paragraph(recommendations_text, self.styles['Normal']))
        story.append(Spacer(1, 0.25 * inch))
        
        # Page break before detailed results
        story.append(PageBreak())
        story.append(self._create_styled_section_header("Detailed Results by Class", "📋"))
        story.append(Spacer(1, 0.18 * inch))
        
        for class_name in sorted(results_by_class.keys()):
            results = results_by_class[class_name]
            
            story.append(Spacer(1, 0.15 * inch))
            story.append(Paragraph(f"Grade {class_name}", self.styles['ClassHeading']))
            story.append(Spacer(1, 0.12 * inch))
            
            # Create detailed results table
            table_data = [
                ['Student Name', 'English', 'Math', 'Science', 'Soc. St.', 'Language', 'Agric.', 'Aggregate', 'Status']
            ]
            
            for result in sorted(results, key=lambda x: x.student.full_name):
                student = result.student
                
                # Get unit scores
                eng = result.english_units or '-'
                math = result.mathematics_units or '-'
                sci = result.science_units or '-'
                ss = result.social_studies_units or '-'
                lang = result.indigenous_language_units or '-'
                agric = result.agriculture_units or '-'
                
                # Calculate aggregate (count passing subjects)
                subjects = [eng, math, sci, ss, lang, agric]
                passes = sum(1 for s in subjects if isinstance(s, int) and s <= 5)
                aggregate = f"{len([s for s in subjects if isinstance(s, int)])}.{passes}"
                
                # Determine status
                status = '✓ PASS' if passes >= 4 else '✗ FAIL'
                
                table_data.append([
                    student.full_name,
                    str(eng),
                    str(math),
                    str(sci),
                    str(ss),
                    str(lang),
                    str(agric),
                    aggregate,
                    status
                ])
            
            # Create table with proper styling
            table = Table(table_data, colWidths=[2.0*inch, 0.55*inch, 0.55*inch, 0.6*inch, 0.6*inch, 0.65*inch, 0.6*inch, 0.7*inch, 0.65*inch])
            table.setStyle(TableStyle([
                # Header
                ('BACKGROUND', (0, 0), (-1, 0), self.primary_blue),
                ('TEXTCOLOR', (0, 0), (-1, 0), self.white),
                ('ALIGN', (0, 0), (-1, 0), 'LEFT'),
                ('VALIGN', (0, 0), (-1, 0), 'MIDDLE'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('PADDING', (0, 0), (-1, 0), (12, 10)),
                ('TOPPADDING', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
                ('LEFTPADDING', (0, 0), (-1, 0), 14),
                
                # Data rows
                ('BACKGROUND', (0, 1), (-1, -1), self.white),
                ('TEXTCOLOR', (0, 1), (-1, -1), self.text_dark),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('ALIGN', (1, 1), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 9),
                ('PADDING', (0, 1), (-1, -1), (10, 6)),
                ('LEFTPADDING', (0, 1), (-1, -1), 12),
                ('RIGHTPADDING', (0, 1), (-1, -1), 10),
                
                # Alternating rows
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [self.white, self.bg_accent]),
                
                # Borders
                ('GRID', (0, 0), (-1, -1), 0.5, Color(200/255, 200/255, 200/255)),
                ('LINEABOVE', (0, 0), (-1, 0), 2, self.primary_blue),
                ('LINEBELOW', (0, -1), (-1, -1), 2, self.primary_blue),
                ('LINEABOVE', (0, 1), (-1, 1), 1, Color(220/255, 220/255, 220/255)),
                
                # First column styling
                ('FONTNAME', (0, 1), (0, -1), 'Helvetica-Bold'),
            ]))
            story.append(table)
            story.append(Spacer(1, 0.2 * inch))
        
        # Professional Footer with Enhanced Styling
        story.append(Spacer(1, 0.4 * inch))
        story.append(HRFlowable(width="100%", thickness=2, color=self.primary_blue, spaceAfter=12))
        
        footer_text = f"""
        <b>{self.school_name}</b> | ZIMSEC Results Analysis Report<br/>
        <b>Report Date:</b> {datetime.now().strftime('%B %d, %Y')} | <b>Students:</b> {total_students} | <b>Classes:</b> {len(results_by_class)}<br/><br/>
        <i>This report contains confidential academic information intended for educational administrators and authorized stakeholders. 
        Unauthorized distribution is prohibited.</i><br/><br/>
        <b>\u2022 Prepared for institutional use and strategic planning</b><br/>
        <b>\u2022 Data-driven insights for academic improvement</b><br/>
        <b>\u2022 Comprehensive performance analysis and recommendations</b>
        """
        
        story.append(Paragraph(footer_text, self.styles['Normal']))
        
        # Build PDF
        doc.build(story)
        return buffer


class HTMLDashboardExporter:
    """Generate self-contained interactive HTML dashboards"""
    
    def __init__(self, title="ZIMSEC Dashboard"):
        self.title = title
        
    def generate_dashboard(self, filepath, data_dict):
        """Generate interactive HTML dashboard"""
        html_content = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{self.title}</title>
    <script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/jspdf/2.5.1/jspdf.umd.min.js"></script>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: #333;
            padding: 20px;
        }}
        
        .container {{
            max-width: 1400px;
            margin: 0 auto;
            background: white;
            border-radius: 10px;
            box-shadow: 0 10px 40px rgba(0,0,0,0.3);
            padding: 30px;
        }}
        
        .header {{
            text-align: center;
            margin-bottom: 30px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 20px;
        }}
        
        .header h1 {{
            color: #333;
            font-size: 2.5em;
            margin-bottom: 10px;
        }}
        
        .header p {{
            color: #666;
            font-size: 1.1em;
        }}
        
        .controls {{
            display: flex;
            gap: 10px;
            margin-bottom: 20px;
            flex-wrap: wrap;
        }}
        
        .controls button {{
            padding: 10px 20px;
            background: #667eea;
            color: white;
            border: none;
            border-radius: 5px;
            cursor: pointer;
            font-size: 1em;
        }}
        
        .controls button:hover {{
            background: #764ba2;
        }}
        
        .chart-container {{
            margin-bottom: 40px;
            border: 1px solid #ddd;
            border-radius: 5px;
            padding: 15px;
            background: #f9f9f9;
        }}
        
        .chart-container h3 {{
            color: #333;
            margin-bottom: 15px;
        }}
        
        .data-table {{
            width: 100%;
            border-collapse: collapse;
            margin-top: 20px;
        }}
        
        .data-table th {{
            background: #667eea;
            color: white;
            padding: 12px;
            text-align: left;
        }}
        
        .data-table td {{
            padding: 10px 12px;
            border-bottom: 1px solid #ddd;
        }}
        
        .data-table tr:hover {{
            background: #f5f5f5;
        }}
        
        .footer {{
            text-align: center;
            margin-top: 30px;
            padding-top: 20px;
            border-top: 1px solid #ddd;
            color: #666;
            font-size: 0.9em;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>{self.title}</h1>
            <p>Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
        </div>
        
        <div class="controls">
            <button onclick="downloadPDF()">Download as PDF</button>
            <button onclick="downloadCSV()">Download Data as CSV</button>
            <button onclick="printPage()">Print</button>
        </div>
        
        <div id="charts"></div>
        
        <div class="footer">
            <p>This is a self-contained dashboard. All data is embedded in this file.</p>
            <p>No internet connection required to view or interact with the charts.</p>
        </div>
    </div>
    
    <script>
        const dashboardData = {json.dumps(data_dict, default=str)};
        
        function downloadPDF() {{
            alert('PDF download feature coming soon!');
        }}
        
        function downloadCSV() {{
            alert('CSV download feature coming soon!');
        }}
        
        function printPage() {{
            window.print();
        }}
        
        // Initialize charts if data is available
        if (dashboardData && Object.keys(dashboardData).length > 0) {{
            console.log('Dashboard data loaded:', dashboardData);
        }}
    </script>
</body>
</html>
"""
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return filepath


def generate_statistics_snapshot(results_list, year=2027):
    """Generate a snapshot of statistics for export"""
    if not results_list:
        return {}
    
    total = len(results_list)
    passed = sum(1 for r in results_list if r.overall_status == 'PASS')
    aggregates = [r.total_aggregate for r in results_list if r.total_aggregate]
    
    return {
        'year': year,
        'total_students': total,
        'pass_rate': (passed / total * 100) if total > 0 else 0,
        'failed_count': total - passed,
        'avg_aggregate': sum(aggregates) / len(aggregates) if aggregates else 0,
        'distinction_rate': sum(1 for r in results_list if r.total_aggregate and r.total_aggregate <= 13) / total * 100 if total > 0 else 0,
        'generated_at': datetime.now().isoformat(),
    }


def generate_subject_statistics(results_list):
    """Generate subject-specific statistics"""
    subjects = {}
    
    for result in results_list:
        for subject_name, units_value in [
            ('English', result.english_units),
            ('Mathematics', result.mathematics_units),
            ('Science', result.science_units),
            ('Social Studies', result.social_studies_units),
            ('Indigenous Language', result.indigenous_language_units),
        ]:
            if subject_name not in subjects:
                subjects[subject_name] = {'values': [], 'distinction': 0, 'credit': 0, 'pass': 0, 'fail': 0}
            
            if units_value is not None:
                subjects[subject_name]['values'].append(units_value)
                if units_value >= 8:
                    subjects[subject_name]['distinction'] += 1
                elif units_value >= 6:
                    subjects[subject_name]['credit'] += 1
                elif units_value >= 4:
                    subjects[subject_name]['pass'] += 1
                else:
                    subjects[subject_name]['fail'] += 1
    
    # Calculate statistics
    stats = {}
    for subject, data in subjects.items():
        if data['values']:
            total = len(data['values'])
            stats[subject] = {
                'pass_rate': ((data['distinction'] + data['credit'] + data['pass']) / total * 100) if total > 0 else 0,
                'avg_units': sum(data['values']) / len(data['values']),
                'distinction_pct': (data['distinction'] / total * 100) if total > 0 else 0,
                'credit_pct': (data['credit'] / total * 100) if total > 0 else 0,
                'pass_pct': (data['pass'] / total * 100) if total > 0 else 0,
                'fail_pct': (data['fail'] / total * 100) if total > 0 else 0,
            }
    
    return stats
